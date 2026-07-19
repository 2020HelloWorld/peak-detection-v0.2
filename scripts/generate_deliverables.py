#!/usr/bin/env python3
"""Generate the v0.5.2 Chinese PPTX and DOCX project deliverables.

All charts, counts, case metrics, and chromatogram images are loaded from the
current detector output. Re-running this script deterministically replaces the
two generated files under ``docs/``.
"""

from __future__ import annotations

import io
import json
import math
import re
import sys
from collections import Counter
from datetime import date
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
VENDOR = ROOT / ".vendor"
if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPTX_PATH = ROOT / "docs" / "色谱峰检测成果展示_v0.5.2_全量41条结果版.pptx"
MD_PATH = ROOT / "docs" / "技术总结_v0.5.2.md"
DOCX_PATH = ROOT / "docs" / "技术总结_v0.5.2.docx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft YaHei"

NAVY = "102A43"
NAVY2 = "243B53"
INK = "243B53"
MUTED = "627D98"
TEAL = "0E8F87"
GREEN = "2E9E5B"
ORANGE = "F3A712"
RED = "D64550"
PURPLE = "6C5CE7"
BLUE = "2F80ED"
WHITE = "FFFFFF"
PALE = "F4F7FA"
PALE_BLUE = "EAF2FF"
PALE_GREEN = "EAF7EF"
PALE_YELLOW = "FFF7DF"
PALE_RED = "FDEDEF"
LINE = "D9E2EC"

plt.rcParams["font.sans-serif"] = [FONT, "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def docx_rgb(value: str):
    """Return python-docx's RGBColor (distinct from python-pptx's class)."""
    from docx.shared import RGBColor as DocxRGBColor

    return DocxRGBColor.from_string(value)


def select_result_root() -> Path:
    candidates = [
        ROOT / "outputs" / "reliable_detector_results",
        ROOT / "outputs" / "reliable_detector_results_v052_final",
        ROOT / "outputs" / "raw_zoom_validation",
        ROOT / "outputs" / "reliable_detector_results_v052",
    ]
    required = (
        "all_detected_features.csv",
        "file_quality_summary.csv",
        "learned_peak_template.json",
        "manifest.json",
    )
    for folder in candidates:
        if not folder.is_dir() or not all((folder / name).exists() for name in required):
            continue
        manifest = json.loads((folder / "manifest.json").read_text(encoding="utf-8"))
        if str(manifest.get("algorithm_version", "")).startswith("0.5.2"):
            return folder
    raise FileNotFoundError("未找到完整的 v0.5.2 检测结果目录")


def select_zoom_root(result_root: Path) -> Path | None:
    for folder in (
        result_root / "raw_zoom_plots",
        ROOT / "outputs" / "raw_zoom_validation" / "raw_zoom_plots",
    ):
        if folder.is_dir() and any(folder.glob("*.png")):
            return folder
    return None


def read_csv(path: Path) -> pd.DataFrame:
    return pd.read_csv(path, encoding="utf-8-sig")


def find_image(folder: Path | None, token: str) -> Path | None:
    if folder is None or not folder.exists():
        return None
    matches = sorted(folder.glob(f"*{token}*.png"))
    return matches[0] if matches else None


def crop_image(path: Path, crop: tuple[float, float, float, float] | None = None) -> io.BytesIO:
    with Image.open(path) as source:
        image = source.convert("RGB")
        if crop:
            left, top, right, bottom = crop
            width, height = image.size
            image = image.crop(
                (int(left * width), int(top * height), int(right * width), int(bottom * height))
            )
        stream = io.BytesIO()
        image.save(stream, "PNG", optimize=True)
    stream.seek(0)
    return stream


def chart_image(fig: plt.Figure) -> io.BytesIO:
    stream = io.BytesIO()
    fig.savefig(stream, format="png", dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    stream.seek(0)
    return stream


def add_run(paragraph, text: str, size: float, color: str, bold: bool = False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return run


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    add_run(paragraph, text, size, color, bold)
    return shape


def add_lines(slide, lines, x: float, y: float, w: float, h: float, *, size=13.5, spacing=6):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.07)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.04)
    for index, (text, color, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.04
        add_run(paragraph, text, size, color, bold)
    return shape


def add_box(slide, x: float, y: float, w: float, h: float, fill=WHITE, line=LINE, rounded=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def add_card(slide, value: str, label: str, x: float, y: float, w: float, color: str, note=""):
    add_box(slide, x, y, w, 1.14)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.14)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(color)
    accent.line.fill.background()
    add_text(slide, value, x + 0.20, y + 0.10, w - 0.30, 0.46, size=25, color=color, bold=True)
    add_text(slide, label, x + 0.20, y + 0.61, w - 0.30, 0.23, size=11.5, bold=True)
    if note:
        add_text(slide, note, x + 0.20, y + 0.88, w - 0.30, 0.17, size=8.3, color=MUTED)


def add_picture(slide, image, x: float, y: float, w: float, h: float, border=True):
    stream = crop_image(image) if isinstance(image, Path) else image
    stream.seek(0)
    with Image.open(stream) as opened:
        iw, ih = opened.size
    stream.seek(0)
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        add_box(slide, x, y, w, h, rounded=False)
    return slide.shapes.add_picture(stream, Inches(px), Inches(py), Inches(pw), Inches(ph))


def new_slide(
    prs: Presentation,
    title: str,
    number: int,
    source: str,
    *,
    title_size: float = 24,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.11), SLIDE_H)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(TEAL)
    stripe.line.fill.background()
    add_text(slide, title, 0.46, 0.24, 11.65, 0.48, size=title_size, color=NAVY, bold=True)
    add_text(slide, f"{number:02d}", 12.15, 0.27, 0.55, 0.28, size=11, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.47), Inches(0.83), Inches(12.15), Inches(0.015)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(LINE)
    rule.line.fill.background()
    add_text(
        slide,
        f"色谱峰检测成果展示 · v0.5.2  |  数据源：{source}",
        0.47,
        7.17,
        11.9,
        0.16,
        size=8.2,
        color=MUTED,
    )
    return slide


def add_case_banner(slide, label: str, text: str, color=TEAL):
    add_box(slide, 0.55, 0.98, 12.08, 0.52, fill=PALE)
    add_text(slide, label, 0.75, 1.10, 1.30, 0.22, size=11.5, color=color, bold=True)
    add_text(slide, text, 1.95, 1.08, 10.35, 0.25, size=12.5, bold=True)


def status_chart(counts: Counter) -> io.BytesIO:
    values = [counts["confirmed"], counts["review"], counts["artifact"]]
    labels = ["确认", "待复核", "干扰/伪峰"]
    colors = [f"#{GREEN}", f"#{ORANGE}", f"#{RED}"]
    fig, ax = plt.subplots(figsize=(5.1, 3.4))
    wedges, _ = ax.pie(
        values,
        colors=colors,
        startangle=90,
        counterclock=False,
        wedgeprops={"width": 0.34, "edgecolor": "white", "linewidth": 2},
    )
    ax.text(0, 0.08, f"{sum(values)}", ha="center", fontsize=26, fontweight="bold", color=f"#{NAVY}")
    ax.text(0, -0.16, "候选事件", ha="center", fontsize=10, color=f"#{MUTED}")
    ax.legend(
        wedges,
        [f"{name} {value}" for name, value in zip(labels, values)],
        loc="lower center",
        bbox_to_anchor=(0.5, -0.18),
        ncol=3,
        frameon=False,
        fontsize=8.5,
    )
    ax.set(aspect="equal")
    fig.tight_layout()
    return chart_image(fig)


def data_chart(summary: pd.DataFrame) -> io.BytesIO:
    names = []
    for folder in summary["folder"].astype(str):
        name = folder.partition("-")[2] or folder
        names.append(name.replace("谱图类型", "").replace("类型", ""))
    counts = Counter(names)
    ordered = sorted(counts.items(), key=lambda pair: pair[1])
    fig, ax = plt.subplots(figsize=(6.2, 3.6))
    bars = ax.barh([item[0] for item in ordered], [item[1] for item in ordered], color=f"#{TEAL}")
    ax.bar_label(bars, padding=4, fontsize=9)
    ax.set_xlabel("独立曲线数")
    ax.grid(axis="x", alpha=0.15)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0)
    fig.tight_layout()
    return chart_image(fig)


def template_chart(template: list[dict]) -> io.BytesIO:
    times = np.asarray([row["retention_time_min"] for row in template], float)
    tolerances = np.asarray([row["rt_tolerance_min"] for row in template], float)
    positions = np.arange(len(template))
    fig, ax = plt.subplots(figsize=(6.8, 3.2))
    ax.errorbar(
        times,
        positions,
        xerr=tolerances,
        fmt="o",
        color=f"#{TEAL}",
        ecolor=f"#{GREEN}",
        elinewidth=8,
        markersize=7,
        alpha=0.85,
    )
    for y, time, tolerance in zip(positions, times, tolerances):
        ax.text(time, y + 0.22, f"{time:.3f} ± {tolerance:.3f}", ha="center", fontsize=8)
    ax.set_yticks(positions, [row["slot"] for row in template])
    ax.set_xlabel("保留时间（min）")
    ax.set_xlim(0, max(times + tolerances) + 0.35)
    ax.grid(axis="x", alpha=0.15)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0)
    fig.tight_layout()
    return chart_image(fig)


def type_chart(features: pd.DataFrame) -> io.BytesIO:
    counts = features.groupby("feature_type_cn", dropna=False).size().nlargest(8).sort_values()
    fig, ax = plt.subplots(figsize=(6.2, 3.9))
    bars = ax.barh(counts.index.astype(str), counts.values, color=f"#{MUTED}")
    ax.bar_label(bars, padding=4, fontsize=8)
    ax.set_xlabel("事件数")
    ax.grid(axis="x", alpha=0.15)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", labelsize=8.2, length=0)
    fig.tight_layout()
    return chart_image(fig)


def confidence_chart(features: pd.DataFrame) -> io.BytesIO:
    keys = ["confirmed", "review", "artifact"]
    names = ["确认", "待复核", "干扰/伪峰"]
    values = [features.loc[features["status"] == key, "peak_confidence_percent"].astype(float) for key in keys]
    fig, ax = plt.subplots(figsize=(5.5, 3.4))
    violin = ax.violinplot(values, showmeans=True, showextrema=False)
    for body, color in zip(violin["bodies"], (GREEN, ORANGE, RED)):
        body.set_facecolor(f"#{color}")
        body.set_edgecolor("none")
        body.set_alpha(0.75)
    violin["cmeans"].set_color(f"#{NAVY}")
    ax.axhline(75, color=f"#{GREEN}", linestyle="--", linewidth=1, label="确认阈值 75%")
    ax.axhline(45, color=f"#{RED}", linestyle=":", linewidth=1, label="伪峰阈值 45%")
    ax.set_xticks([1, 2, 3], names)
    ax.set_ylim(0, 102)
    ax.set_ylabel("峰置信度（%）")
    ax.grid(axis="y", alpha=0.15)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, fontsize=8)
    fig.tight_layout()
    return chart_image(fig)


def weight_bar_chart(
    items: list[tuple[str, int]],
    title: str,
    color: str,
    *,
    label_size: float = 9.0,
) -> io.BytesIO:
    labels = [item[0] for item in items][::-1]
    values = [item[1] for item in items][::-1]
    height = max(3.1, 0.38 * len(items) + 1.0)
    fig, ax = plt.subplots(figsize=(5.8, height))
    bars = ax.barh(labels, values, color=f"#{color}", height=0.62)
    ax.bar_label(bars, labels=[f"{value}%" for value in values], padding=4, fontsize=8.8, fontweight="bold")
    ax.set_xlim(0, max(values) + 8)
    ax.set_xlabel("权重")
    ax.set_title(title, loc="left", fontsize=13, fontweight="bold", color=f"#{NAVY}")
    ax.grid(axis="x", alpha=0.13)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0, labelsize=label_size)
    fig.tight_layout()
    return chart_image(fig)


def result_mix_by_workload_chart(features: pd.DataFrame) -> io.BytesIO:
    """Compare final output composition across the five supplied workloads.

    This is deliberately an output-mix chart, not an accuracy chart: there is
    no row-level human truth table yet.  Percentages therefore use detected
    candidates as the denominator within each workload.
    """
    folder_names = {
        "1-基准谱图类型": "参考标气",
        "2-鼓包谱图类型": "鼓包",
        "3-噪音谱图类型": "噪声",
        "4-负峰类型": "负峰",
        "5-电信号干扰类型": "电信号干扰",
    }
    order = list(folder_names.values())
    table = pd.crosstab(
        features["folder"].map(folder_names).fillna(features["folder"].astype(str)),
        features["status"],
    ).reindex(index=order, columns=["confirmed", "review", "artifact"], fill_value=0)
    totals = table.sum(axis=1)
    percent = table.div(totals.replace(0, np.nan), axis=0).fillna(0) * 100

    fig, ax = plt.subplots(figsize=(7.7, 3.8))
    left = np.zeros(len(table))
    colors = [f"#{GREEN}", f"#{ORANGE}", f"#{RED}"]
    labels = ["确认", "待复核", "干扰/伪峰"]
    for key, label, color in zip(table.columns, labels, colors):
        values = percent[key].to_numpy(float)
        bars = ax.barh(table.index, values, left=left, color=color, label=label, height=0.62)
        for bar, value, count in zip(bars, values, table[key].to_numpy(int)):
            if value >= 8:
                ax.text(
                    bar.get_x() + bar.get_width() / 2,
                    bar.get_y() + bar.get_height() / 2,
                    f"{value:.0f}%\n({count})",
                    ha="center",
                    va="center",
                    fontsize=8,
                    color="white" if key != "review" else f"#{NAVY}",
                    fontweight="bold",
                )
        left += values
    for y, total in enumerate(totals.to_numpy(int)):
        ax.text(101.5, y, f"n={total}", va="center", fontsize=8.5, color=f"#{MUTED}")
    ax.set_xlim(0, 112)
    ax.set_xlabel("该工况全部候选事件中的状态占比")
    ax.invert_yaxis()
    ax.grid(axis="x", alpha=0.12)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0)
    ax.legend(loc="lower center", bbox_to_anchor=(0.46, -0.30), ncol=3, frameon=False)
    fig.tight_layout()
    return chart_image(fig)


def add_comparison_headers(slide) -> None:
    add_text(slide, "案例", 0.55, 1.00, 0.72, 0.24, size=10.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "原始数据局部放大（未预处理）", 1.42, 1.00, 4.38, 0.24, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 5.83, 1.00, 0.30, 0.24, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "最终检测与分类标注", 6.16, 1.00, 4.38, 0.24, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "对比结论", 10.68, 1.00, 1.90, 0.24, size=10.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)


def add_comparison_row(
    slide,
    *,
    y: float,
    case_label: str,
    raw_image: Path | None,
    raw_crop: tuple[float, float, float, float],
    detected_image: Path | None,
    detected_crop: tuple[float, float, float, float],
    conclusion: list[tuple[str, str, bool]],
    color: str,
) -> None:
    row_h = 2.48
    add_box(slide, 0.55, y, 0.72, row_h, fill=color, line=color)
    add_text(
        slide,
        case_label,
        0.58,
        y + 0.74,
        0.66,
        0.80,
        size=15,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if raw_image:
        add_picture(slide, crop_image(raw_image, raw_crop), 1.42, y, 4.38, row_h)
    else:
        add_box(slide, 1.42, y, 4.38, row_h, fill=PALE)
        add_text(slide, "未找到原始分段图", 1.72, y + 1.02, 3.78, 0.28, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 5.83, y + 1.00, 0.30, 0.32, size=21, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    if detected_image:
        add_picture(slide, crop_image(detected_image, detected_crop), 6.16, y, 4.38, row_h)
    else:
        add_box(slide, 6.16, y, 4.38, row_h, fill=PALE)
        add_text(slide, "未找到最终检测图", 6.46, y + 1.02, 3.78, 0.28, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_box(slide, 10.68, y, 1.90, row_h, fill=PALE, line=color)
    add_lines(slide, conclusion, 10.83, y + 0.20, 1.60, row_h - 0.36, size=9.4, spacing=4)


def event_near(features: pd.DataFrame, prefix: str, time_min: float) -> pd.Series:
    subset = features[features["file"].astype(str).str.startswith(prefix)].copy()
    if subset.empty:
        raise KeyError(prefix)
    index = (subset["apex_time_min"].astype(float) - time_min).abs().idxmin()
    return subset.loc[index]


def file_near(summary: pd.DataFrame, prefix: str) -> pd.Series:
    subset = summary[summary["file"].astype(str).str.startswith(prefix)]
    if subset.empty:
        raise KeyError(prefix)
    return subset.iloc[0]


def build_presentation(result_root: Path, zoom_root: Path | None) -> Presentation:
    features = read_csv(result_root / "all_detected_features.csv")
    summary = read_csv(result_root / "file_quality_summary.csv")
    template = json.loads((result_root / "learned_peak_template.json").read_text(encoding="utf-8"))
    manifest = json.loads((result_root / "manifest.json").read_text(encoding="utf-8"))
    counts = Counter(features["status"].astype(str))
    expected = {"confirmed": 231, "review": 380, "artifact": 49}
    found = {key: counts[key] for key in expected}
    if found != expected:
        raise RuntimeError(f"正式统计不一致：found={found}, expected={expected}")

    plots = result_root / "plots"
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    source = result_root.name

    # 01 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-1.0), Inches(5.2), Inches(5.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(TEAL)
    circle.fill.transparency = 18
    circle.line.fill.background()
    add_text(slide, "色谱峰检测成果展示", 0.82, 1.28, 8.8, 0.75, size=34, color=WHITE, bold=True)
    add_text(slide, "从原始谱图到峰位、峰型、边界与置信度输出", 0.84, 2.16, 9.0, 0.45, size=20, color="D9EAF2")
    add_text(slide, "v0.5.2", 0.84, 3.02, 1.25, 0.34, size=18, color=WHITE, bold=True)
    add_text(slide, "41 条独立曲线 · 6 个参考模板槽 · 可审计 CSV + 图像", 2.05, 3.05, 7.8, 0.28, size=13, color="B7D9E8")
    add_box(slide, 0.84, 4.26, 7.1, 1.32, fill=NAVY2, line="3E5C76")
    add_lines(
        slide,
        [
            ("聚焦已实现能力、代表性修正、可验证成果与边界", WHITE, True),
            ("所有数字和案例均直接读取当前 v0.5.2 正式输出", "B7D9E8", False),
        ],
        1.06,
        4.56,
        6.6,
        0.75,
        size=14.5,
        spacing=8,
    )
    add_text(slide, str(date.today()), 0.85, 6.74, 2.0, 0.22, size=9.5, color="8FB8CC")
    add_text(slide, "01", 12.05, 6.70, 0.55, 0.26, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    # 02 Outcome
    slide = new_slide(prs, "一页看懂当前成果", 2, source)
    add_picture(slide, status_chart(counts), 0.55, 1.07, 5.15, 4.85)
    add_card(slide, str(counts["confirmed"]), "已确认峰", 6.02, 1.20, 2.05, GREEN, "峰置信度 ≥ 75%")
    add_card(slide, str(counts["review"]), "待复核事件", 8.27, 1.20, 2.05, ORANGE, "保留完整证据")
    add_card(slide, str(counts["artifact"]), "干扰／伪峰", 10.52, 1.20, 2.05, RED, "电形态或低可信")
    add_box(slide, 6.02, 2.66, 6.55, 3.22, fill=PALE)
    add_lines(
        slide,
        [
            ("✓ 保留未经处理的原始谱线，便于逐图对照", GREEN, True),
            ("✓ 峰顶时间、校正值、起止边界与综合置信度同步输出", GREEN, True),
            ("✓ 正峰、负峰、峰顶尖点、电干扰、鼓包与背景残差可区分", GREEN, True),
            ("✓ review 不隐藏、不强行确认，可直接形成复核清单", GREEN, True),
            ("已形成可批量运行、可解释、可回归验证的工程初版。", NAVY, True),
        ],
        6.29,
        2.96,
        6.0,
        2.55,
        size=13.5,
        spacing=10,
    )

    # 03 Dataset
    slide = new_slide(prs, "项目目标与数据覆盖", 3, source)
    add_picture(slide, data_chart(summary), 0.55, 1.08, 6.45, 4.90)
    add_card(slide, str(manifest.get("independent_curves", len(summary))), "独立曲线", 7.28, 1.18, 1.55, BLUE)
    add_card(slide, str(manifest.get("reference_curves", 5)), "参考标气", 8.99, 1.18, 1.55, TEAL)
    add_card(slide, str(manifest.get("learned_template_slots", len(template))), "模板槽", 10.70, 1.18, 1.55, PURPLE)
    add_box(slide, 7.28, 2.68, 5.00, 3.28, fill=PALE_BLUE, line="C9DDF7")
    add_lines(
        slide,
        [
            ("目标", BLUE, True),
            ("定位每个候选事件，判断峰型，并给出可核验的置信度与边界。", INK, False),
            ("覆盖", BLUE, True),
            ("参考标气、鼓包、噪声、负峰、电信号干扰五类工况。", INK, False),
            ("评估原则", BLUE, True),
            ("算法输出是待验证结论；没有人工标签时不宣称等同真值。", RED, True),
        ],
        7.55,
        2.94,
        4.45,
        2.65,
        size=13.7,
        spacing=6,
    )

    # 04 Complete legend
    slide = new_slide(prs, "检测结果图完整图例｜曲线、区域、状态与形态符号", 4, source)
    a7_plot = find_image(plots, "A7-")
    add_box(slide, 0.55, 1.05, 4.03, 5.85, fill=PALE_BLUE, line="C9DDF7")
    add_text(slide, "三层曲线", 0.82, 1.28, 2.2, 0.32, size=18, color=BLUE, bold=True)
    add_lines(
        slide,
        [
            ("第一层｜raw signal", NAVY, True),
            ("━━ 深灰 raw samples：直接连接原始 CSV 采样，不平滑、不扣基线。", "555555", False),
            ("第二层｜signal / baseline", NAVY, True),
            ("━━ 浅灰 raw reference：原始谱线的位置参照。", "8C8C8C", False),
            ("━━ 红色 robust main baseline：正峰检测使用的稳健主基线。", RED, False),
            ("━━ 紫色 rolling-ball local background：鼓包和局部背景模型。", PURPLE, False),
            ("第三层｜corrected signal", NAVY, True),
            ("━━ 蓝色 preprocessed positive：主基线扣除后的正峰检测分支。", BLUE, False),
            ("━━ 灰色 signed branch：保留正负号，供负峰搜索与确认。", MUTED, False),
            ("━━ 黑灰水平线：基线校正后的 0 参考线。", INK, False),
        ],
        0.80,
        1.74,
        3.55,
        4.86,
        size=10.4,
        spacing=5,
    )

    add_box(slide, 4.72, 1.05, 3.68, 5.85, fill=PALE_YELLOW, line="F0D78A")
    add_text(slide, "区域、边界与标签", 4.99, 1.28, 2.9, 0.32, size=18, color=ORANGE, bold=True)
    add_lines(
        slide,
        [
            ("▧ 淡黄色区域", ORANGE, True),
            ("宽背景／鼓包诊断区；其中仍允许存在证据充分的真实峰。", INK, False),
            ("▥ 浅绿色竖向矩形", GREEN, True),
            ("T1–T6 保留时间容差窗；只是模板位置，不是峰或峰边界。", INK, False),
            ("▷ 青绿空心向右三角", TEAL, True),
            ("仅用于 confirmed 峰，表示算法起点。", INK, False),
            ("◁ 青绿空心向左三角", TEAL, True),
            ("仅用于 confirmed 峰，表示算法终点；两个三角都朝向峰区间内部。", INK, False),
            ("标签与细引导线", NAVY, True),
            ("显示峰型、置信度、峰顶时间／校正值及确认峰起止时间；引导线只连接文字与对应峰顶。", INK, False),
            ("边界为自动近基线边界，不等同于人工确认的最终积分边界。", RED, True),
        ],
        4.97,
        1.74,
        3.18,
        4.88,
        size=10.2,
        spacing=4,
    )

    add_box(slide, 8.54, 1.05, 4.04, 5.85, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "状态与形态叠加", 8.81, 1.28, 2.9, 0.32, size=18, color=GREEN, bold=True)
    add_lines(
        slide,
        [
            ("● 绿色圆点 confirmed：置信度达到确认线。", GREEN, True),
            ("▲ 橙色三角 review：有峰证据，但需人工复核。", ORANGE, True),
            ("× 红色叉号 artifact / interference：更像电干扰或非色谱伪峰。", RED, True),
            ("◆ 灰色菱形 likely noise：低置信度峰／噪声候选。", MUTED, True),
            ("· 灰色小点 structured background residual：结构化背景残差，保留审计但不确认。", "777777", True),
            ("□ 紫色空心方框 gentle broad peak candidate：平缓宽峰形态叠加。", PURPLE, True),
            ("P 紫色空心加号形 chromatographic peak with apex spike：峰体真实但峰顶带尖点。", PURPLE, True),
            ("状态符号回答“可靠性等级”；紫色叠加符号回答“峰的形态”。同一事件可以同时出现两类符号。", NAVY, True),
        ],
        8.79,
        1.74,
        3.54,
        4.86,
        size=10.3,
        spacing=5,
    )

    # 05 Template
    slide = new_slide(prs, "参考模板：T1–T6 是相对位置槽，不是化学名称", 5, source)
    add_picture(slide, template_chart(template), 0.55, 1.15, 7.22, 4.62)
    add_box(slide, 8.02, 1.15, 4.58, 4.62, fill=PALE_GREEN, line="C8E8D2")
    template_lines = [("模板来自 5 份参考标气的稳定峰族", GREEN, True)]
    template_lines.extend(
        (
            f"{row['slot']}  {row['retention_time_min']:.3f} ± {row['rt_tolerance_min']:.3f} min · 支持 {row['support']}/5",
            INK,
            False,
        )
        for row in template
    )
    template_lines.extend(
        [
            ("跨机型含义", RED, True),
            ("当前仅表示相似保留时间族；补充通道／组分映射后才能命名 H₂、CH₄ 等。", INK, False),
        ]
    )
    add_lines(slide, template_lines, 8.27, 1.40, 4.10, 4.15, size=11.8, spacing=5)
    add_text(slide, "模板提供匹配证据；形态异常仍可判为电干扰，保留时间不能“救回”强毛刺。", 0.70, 6.07, 11.7, 0.40, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 06 Confidence overview
    slide = new_slide(prs, "置信度判别标准｜分数用途、状态阈值与动态修正", 6, source)
    add_box(slide, 0.55, 1.00, 12.03, 0.58, fill=PALE, line=LINE)
    add_text(
        slide,
        "peak_confidence = 各项 0～1 证据分 × 固定权重求和，再乘工况修正系数，最后截断到 0～1。",
        0.78,
        1.16,
        11.58,
        0.25,
        size=14,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_box(slide, 0.55, 1.80, 3.85, 4.98, fill=PALE_BLUE, line="C9DDF7")
    add_text(slide, "峰存在置信度", 0.82, 2.05, 2.5, 0.33, size=19, color=BLUE, bold=True)
    add_lines(
        slide,
        [
            ("peak_confidence 回答：该事件像不像真实色谱峰。", NAVY, True),
            ("SNR：约 3～30 采用对数映射。", INK, False),
            ("相对突出度：0.0008～0.03 采用对数映射；高噪声另设 0.01 保护门。", INK, False),
            ("峰宽：与全部参考模板宽度比较，取最佳相似度。", INK, False),
            ("对称性：由 0.20～0.90 映射到 0～1。", INK, False),
            ("非平顶性：由 top_width_ratio 映射，顶部越异常越低。", INK, False),
            ("负峰另外使用带符号下探深度和原始双侧谷深。", PURPLE, True),
            ("它是可解释的概率样评分，尚不是经人工真值校准的统计概率。", RED, True),
        ],
        0.80,
        2.52,
        3.35,
        3.95,
        size=10.9,
        spacing=6,
    )

    add_box(slide, 4.62, 1.80, 3.20, 4.98, fill=WHITE, line=LINE)
    add_text(slide, "状态阈值", 4.90, 2.05, 2.0, 0.33, size=19, color=NAVY, bold=True)
    add_box(slide, 4.91, 2.55, 2.62, 0.92, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "≥ 75%", 5.13, 2.66, 1.10, 0.34, size=22, color=GREEN, bold=True)
    add_text(slide, "confirmed", 6.15, 2.71, 1.14, 0.24, size=11.5, color=GREEN, bold=True)
    add_box(slide, 4.91, 3.68, 2.62, 1.18, fill=PALE_RED, line="F4BCC2")
    add_text(slide, "< 45%", 5.13, 3.80, 1.10, 0.34, size=22, color=RED, bold=True)
    add_text(slide, "artifact", 6.15, 3.85, 1.14, 0.24, size=11.5, color=RED, bold=True)
    add_text(slide, "仅限电事件或 uncertain_peak_or_noise", 5.12, 4.33, 2.10, 0.27, size=8.8, color=RED, bold=True)
    add_box(slide, 4.91, 5.08, 2.62, 0.92, fill=PALE_YELLOW, line="F0D78A")
    add_text(slide, "其余", 5.13, 5.20, 1.10, 0.34, size=22, color=ORANGE, bold=True)
    add_text(slide, "review", 6.15, 5.25, 1.14, 0.24, size=11.5, color=ORANGE, bold=True)
    add_text(slide, "低于 45% 并不自动等于 artifact；类型条件必须同时成立。", 4.94, 6.22, 2.54, 0.36, size=9.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 8.04, 1.80, 4.54, 4.98, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "独立模板分与主要修正", 8.31, 2.05, 3.3, 0.33, size=19, color=GREEN, bold=True)
    add_lines(
        slide,
        [
            ("template_confidence（不替代峰存在分）", NAVY, True),
            ("保留时间 50%｜宽度 30%｜对称性 15%｜非平顶性 5%", GREEN, True),
            ("主要乘法修正", NAVY, True),
            ("×0.75 高噪声且突出度不足；×0.65 靠近采集边缘", INK, False),
            ("×0.25 明显纯电尖峰；×0.55 峰顶异常但证据矛盾", INK, False),
            ("×0.90 真实峰体带峰顶尖点；×0.80 超宽鼓包／基线", INK, False),
            ("×0.65 负峰与脉冲重合；×0.65 两侧正峰夹出的浅谷", INK, False),
            ("结构化背景残差：峰宽奖励减半，置信度上限 0.74。", ORANGE, True),
            ("模板命中只能提供匹配证据，不能覆盖强电形态。", RED, True),
        ],
        8.30,
        2.52,
        4.00,
        3.98,
        size=10.7,
        spacing=5,
    )

    # 07 Primary confidence weights
    slide = new_slide(prs, "峰存在置信度权重｜正峰与负峰使用不同证据组合", 7, source)
    positive_weights = [
        ("局部／有效 SNR", 25),
        ("相对突出度", 25),
        ("参考峰宽相似度", 20),
        ("对称性", 15),
        ("峰顶非平坦度", 15),
    ]
    negative_weights = [
        ("局部 SNR", 18),
        ("相对突出度", 15),
        ("参考峰宽相似度", 12),
        ("对称性", 10),
        ("带符号下探深度", 20),
        ("原始信号双侧谷深", 25),
    ]
    add_picture(slide, weight_bar_chart(positive_weights, "正峰基础权重（合计 100%）", BLUE), 0.55, 1.10, 5.82, 4.85)
    add_picture(slide, weight_bar_chart(negative_weights, "负峰基础权重（合计 100%）", PURPLE), 6.75, 1.10, 5.82, 4.85)
    add_box(slide, 0.72, 6.10, 11.70, 0.72, fill=PALE, line=LINE)
    add_text(
        slide,
        "各证据先独立归一化到 0～1，再按上述权重相加；工况修正系数在加权之后施加。负峰还必须先通过双侧谷深和带符号下探门，未通过者不会进入评分输出。",
        0.98,
        6.30,
        11.18,
        0.30,
        size=11.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 08 Electrical morphology weights
    slide = new_slide(prs, "正候选形态权重｜真实峰体分与电干扰分并行计算", 8, source)
    peak_body_weights = [
        ("FWHM／模板相对宽度", 14),
        ("升降斜率平衡", 10),
        ("对称性", 12),
        ("面积形状", 16),
        ("非平顶性", 8),
        ("局部基线稳定性", 8),
        ("升降连续性", 12),
        ("峰顶面积分散度", 8),
        ("去尖后峰体保留", 8),
        ("二次拟合峰顶超额较小", 4),
    ]
    electrical_weights = [
        ("物理时间 FWHM 过窄", 18),
        ("峰顶面积过度集中", 14),
        ("归一化斜率跳变", 14),
        ("二次拟合峰顶超额", 10),
        ("文件级脉冲密度", 18),
        ("低对称性", 6),
        ("平顶异常", 6),
        ("升降斜率失衡", 8),
        ("去尖后峰体损失", 6),
    ]
    add_picture(slide, weight_bar_chart(peak_body_weights, "peak_body_score｜越高越像完整色谱峰", GREEN, label_size=7.8), 0.45, 1.00, 6.15, 4.95)
    add_picture(slide, weight_bar_chart(electrical_weights, "electrical_interference_score｜越高越像毛刺", RED, label_size=7.8), 6.73, 1.00, 6.15, 4.95)
    add_box(slide, 0.58, 6.05, 12.00, 0.80, fill=PALE_YELLOW, line="F0D78A")
    add_text(
        slide,
        "路由：可靠峰顶异常 + 峰体联合门通过 + peak_body_score ≥ 0.72 且不低于电干扰分 → 带峰顶尖点的色谱峰（×0.90）；电干扰分 ≥ 0.45、峰体门失败且存在尖点或 FWHM < 0.018 min → 纯电尖峰（×0.25）；证据矛盾 → review 倾向（×0.55）。仅正候选计算这两项分数。",
        0.84,
        6.23,
        11.50,
        0.42,
        size=10.4,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 09 Surviving peak-body gate
    slide = new_slide(prs, "带峰顶尖点的色谱峰｜去尖后峰体必须通过的联合门", 9, source)
    add_box(slide, 0.55, 1.02, 12.03, 0.55, fill=PALE, line=LINE)
    add_text(
        slide,
        "该联合门只用于正候选：先证明尖点下面仍有完整峰体，再决定“保留峰”还是“排除毛刺”。",
        0.80,
        1.17,
        11.55,
        0.25,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_box(slide, 0.55, 1.78, 3.72, 5.05, fill=PALE_BLUE, line="C9DDF7")
    add_text(slide, "二选一宽度入口", 0.82, 2.05, 2.7, 0.33, size=19, color=BLUE, bold=True)
    add_box(slide, 0.86, 2.62, 3.10, 1.34, fill=WHITE, line=BLUE)
    add_text(slide, "入口 A｜绝对峰体", 1.08, 2.80, 2.55, 0.28, size=14, color=BLUE, bold=True)
    add_text(slide, "FWHM ≥ 0.060 min\n且 area_shape_factor ≥ 1.15", 1.08, 3.18, 2.56, 0.56, size=12.0, color=INK, bold=True)
    add_text(slide, "或", 2.04, 4.11, 0.50, 0.30, size=17, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 0.86, 4.55, 3.10, 1.72, fill=WHITE, line=GREEN)
    add_text(slide, "入口 B｜模板相对峰体", 1.08, 4.73, 2.55, 0.28, size=14, color=GREEN, bold=True)
    add_text(
        slide,
        "存在模板，width_min 位于\n[0.75×模板宽度下限，1.25×模板宽度上限]\n且 area_shape_factor ≥ 0.85",
        1.08,
        5.11,
        2.58,
        0.92,
        size=10.8,
        color=INK,
        bold=True,
    )

    add_box(slide, 4.49, 1.78, 4.03, 5.05, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "其余条件必须全部满足", 4.76, 2.05, 3.3, 0.33, size=19, color=GREEN, bold=True)
    add_lines(
        slide,
        [
            ("symmetry ≥ 0.75", INK, False),
            ("top_width_ratio ≤ 0.45", INK, False),
            ("baseline_change_ratio ≤ 0.25", INK, False),
            ("min(rise_continuity, fall_continuity) ≥ 0.80", INK, False),
            ("apex_core_area_fraction ≤ 0.35", INK, False),
            ("despiked_height_retention ≥ 0.65", INK, False),
            ("despiked_prominence_retention ≥ 0.65", INK, False),
            ("despiked_area_retention ≥ 0.65", INK, False),
            ("apex_quadratic_excess_ratio ≤ 0.25", INK, False),
            ("上述条件任意一项失败，都不能声称去尖后峰体已经可靠存活。", RED, True),
        ],
        4.77,
        2.60,
        3.50,
        3.95,
        size=10.6,
        spacing=5,
    )

    add_box(slide, 8.74, 1.78, 3.84, 5.05, fill=PALE_YELLOW, line="F0D78A")
    add_text(slide, "最终路由", 9.01, 2.05, 2.4, 0.33, size=19, color=ORANGE, bold=True)
    add_lines(
        slide,
        [
            ("1｜先有可靠尖点证据", NAVY, True),
            ("apex_spike_flag = true。", INK, False),
            ("2｜保留为真实峰", GREEN, True),
            ("联合门通过、peak_body_score ≥ 0.72，且 peak_body_score ≥ electrical_interference_score → peak_with_apex_spike，置信度 ×0.90。", INK, False),
            ("3｜判为纯电尖峰", RED, True),
            ("联合门失败，电干扰分 ≥ 0.45，且有尖点或 FWHM < 0.018 min → 置信度 ×0.25。", INK, False),
            ("4｜保留歧义", ORANGE, True),
            ("两边条件都不足 → electrical_interference_candidate，置信度 ×0.55，通常进入 review。", INK, False),
            ("A7 两个事件通过了此联合门。", GREEN, True),
        ],
        9.00,
        2.58,
        3.30,
        3.98,
        size=10.2,
        spacing=4,
    )

    # 10 Positive and broad/overlap type rules
    slide = new_slide(prs, "峰类型判别标准 I｜正峰、宽峰、鼓包峰与重叠关系", 10, source)
    add_box(slide, 0.55, 1.05, 3.82, 5.82, fill=PALE_BLUE, line="C9DDF7")
    add_text(slide, "模板窗口内", 0.82, 1.30, 2.4, 0.33, size=19, color=BLUE, bold=True)
    add_lines(
        slide,
        [
            ("普通正峰 normal_positive_peak", NAVY, True),
            ("width_min 位于模板宽度上下限，且 symmetry 不低于模板下限。", INK, False),
            ("鼓包上的正峰 positive_peak_on_hump", GREEN, True),
            ("普通峰条件成立，同时 baseline_change_ratio ≥ 0.25，或峰顶／至少 50% 峰窗落入 bump 区。", INK, False),
            ("窄正峰 narrow_positive_peak", NAVY, True),
            ("宽度低于模板下限；若 top_width_ratio ≥ 0.55 或对称性不足，则改为电干扰候选并 ×0.45。", INK, False),
            ("宽峰／重叠峰 broad_or_overlapped_peak", ORANGE, True),
            ("宽度高于模板宽度上限。", INK, False),
            ("其余无法稳定解释的模板内事件 → uncertain_peak_or_noise。", RED, True),
        ],
        0.80,
        1.78,
        3.35,
        4.72,
        size=10.2,
        spacing=4,
    )

    add_box(slide, 4.58, 1.05, 4.22, 5.82, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "模板窗口外", 4.85, 1.30, 2.4, 0.33, size=19, color=GREEN, bold=True)
    add_lines(
        slide,
        [
            ("未映射正峰 unassigned_positive_peak", NAVY, True),
            ("最佳参考宽度得分 ≥ 0.45 且 symmetry ≥ 0.20；表示峰存在但尚未绑定 T1–T6／组分。", INK, False),
            ("平缓宽峰 gentle_broad_peak_candidate", PURPLE, True),
            ("width > 0.80、FWHM ≥ 0.30 min、width/FWHM ≤ 6、symmetry ≥ 0.25、top_width_ratio ≤ 0.65、相对突出度 ≥ 0.003、SNR ≥ 8、baseline_change_ratio ≤ 3、双侧深度 ≥ 0.002、无脉冲且远离边缘。", INK, False),
            ("宽鼓包／基线 broad_hump_or_baseline", ORANGE, True),
            ("width > 0.80 min 但不满足平缓宽峰联合门；置信度 ×0.80。若最终达到确认线，则改名 broad_positive_peak。", INK, False),
            ("模板外极小波纹", MUTED, True),
            ("低于普通 0.0008／高噪声 0.01 相对突出度门且非脉冲、非超宽事件时直接过滤。", INK, False),
        ],
        4.83,
        1.78,
        3.72,
        4.76,
        size=9.7,
        spacing=4,
    )

    add_box(slide, 9.00, 1.05, 3.58, 5.82, fill=PALE_YELLOW, line="F0D78A")
    add_text(slide, "重叠与状态", 9.27, 1.30, 2.4, 0.33, size=19, color=ORANGE, bold=True)
    add_lines(
        slide,
        [
            ("同一模板多候选", NAVY, True),
            ("允许多个 confirmed；模板匹配最好者保留原类型，其余为 secondary_or_overlapping_candidate。", INK, False),
            ("相邻确认峰窗口交叠", NAVY, True),
            ("两个 95% 自动峰窗口相交时标为 overlapping_positive_peak；带峰顶尖点类型优先保留，重叠写入 reasons。", INK, False),
            ("类型与状态分离", RED, True),
            ("“宽峰”“重叠”“未映射”描述形态或模板关系，不自动等于 review。最终状态仍由 peak_confidence 与类型条件共同决定。", INK, False),
            ("当前只识别重叠关系，尚未做多峰参数化解卷积和面积分配。", RED, True),
        ],
        9.25,
        1.78,
        3.05,
        4.75,
        size=10.2,
        spacing=6,
    )

    # 11 Negative, electrical and background-residual rules
    slide = new_slide(prs, "峰类型判别标准 II｜负峰、电干扰与结构化背景残差", 11, source)
    add_box(slide, 0.55, 1.05, 3.78, 5.82, fill="F3EEFF", line="D9CEF7")
    add_text(slide, "负峰分支", 0.82, 1.30, 2.2, 0.33, size=19, color=PURPLE, bold=True)
    add_lines(
        slide,
        [
            ("先通过双门", NAVY, True),
            ("普通工况：bilateral_depth_relative ≥ 0.005 且 signed_depth_relative ≥ 0.008。", INK, False),
            ("高噪声：门槛提高到 0.01 和 0.02。", INK, False),
            ("negative_electrical_spike", RED, True),
            ("width_min < 0.018 min，置信度 ×0.25。", INK, False),
            ("negative_peak / broad_negative_peak", PURPLE, True),
            ("0.018～0.80 min 为负峰；> 0.80 min 为宽负峰。与暂定脉冲重合时 ×0.65。", INK, False),
            ("interpeak_valley_or_negative_peak", ORANGE, True),
            ("signed_depth_relative < 0.05，且前后各 0.8 min 内均存在正峰时视为谷底／负峰歧义，置信度再 ×0.65。", INK, False),
        ],
        0.80,
        1.78,
        3.30,
        4.72,
        size=10.1,
        spacing=4,
    )

    add_box(slide, 4.53, 1.05, 3.82, 5.82, fill=PALE_RED, line="F4BCC2")
    add_text(slide, "峰顶尖点／电干扰", 4.80, 1.30, 2.9, 0.33, size=19, color=RED, bold=True)
    add_lines(
        slide,
        [
            ("可靠峰顶异常入口", NAVY, True),
            ("暂定脉冲重合且二次拟合无法解释的峰顶超额／突出度 ≥ 0.015。", INK, False),
            ("真实峰带峰顶尖点", GREEN, True),
            ("峰体联合门通过、peak_body_score ≥ 0.72 且不低于电干扰分 → peak_with_apex_spike，×0.90。", INK, False),
            ("纯电尖峰／电干扰候选", RED, True),
            ("electrical_interference_score ≥ 0.45、峰体门失败，且存在可靠尖点或 FWHM < 0.018 min → ×0.25。模板外叫 electrical_spike，模板内叫 electrical_interference_candidate。", INK, False),
            ("证据矛盾", ORANGE, True),
            ("有尖点但既不能救回峰体、又不足以判纯电尖峰 → electrical_interference_candidate，×0.55，通常 review。", INK, False),
            ("保留时间命中不能覆盖强电形态。", RED, True),
        ],
        4.78,
        1.78,
        3.32,
        4.75,
        size=9.9,
        spacing=4,
    )

    add_box(slide, 8.55, 1.05, 4.03, 5.82, fill=PALE_YELLOW, line="F0D78A")
    add_text(slide, "结构化背景残差", 8.82, 1.30, 2.9, 0.33, size=19, color=ORANGE, bold=True)
    add_lines(
        slide,
        [
            ("文件级条件满足其一", NAVY, True),
            ("baseline_excursion_ratio ≥ 0.08；或 baseline_curvature ≥ 0.04；或 |endpoint_drift_ratio| ≥ 0.05。", INK, False),
            ("候选级条件必须全部满足", NAVY, True),
            ("relative_prominence < 0.005；width_min < 0.80 min；前后 ±0.50 min 弱候选数 ≥ 2；bilateral_depth_relative < 0.0015。", INK, False),
            ("处理结果", ORANGE, True),
            ("电事件／artifact 优先，不会被残差类型覆盖；其余通过门控者改为 structured_background_residual。", INK, False),
            ("参考峰宽奖励减半，置信度上限 0.74，保持 review 并写入审计 CSV。", RED, True),
            ("它表示更像非平稳背景波纹，不表示原始数据被删除。", NAVY, True),
        ],
        8.80,
        1.78,
        3.52,
        4.75,
        size=10.2,
        spacing=5,
    )

    # 12 A7
    slide = new_slide(prs, "代表性修正｜A7：尖顶不再等于整峰是电干扰", 12, source)
    a7_1 = event_near(features, "A7-", 0.433333333)
    a7_2 = event_near(features, "A7-", 0.565)
    add_case_banner(slide, "修正结果", "两个事件保留“峰顶尖点”标记，同时确认其完整色谱峰体。", GREEN)
    a7_zoom = find_image(zoom_root, "A7-")
    if a7_zoom:
        add_picture(slide, crop_image(a7_zoom, (0.02, 0.04, 0.98, 0.36)), 0.55, 1.68, 6.06, 3.72)
    if a7_plot:
        add_picture(slide, crop_image(a7_plot, (0.03, 0.25, 0.98, 0.99)), 6.80, 1.68, 5.82, 3.72)
    add_card(slide, f"{float(a7_1['peak_confidence_percent']):.1f}%", "0.4333 min 峰置信度", 0.70, 5.62, 2.25, GREEN, f"峰体分 {float(a7_1['peak_body_score']):.3f}")
    add_card(slide, f"{float(a7_2['peak_confidence_percent']):.1f}%", "0.5650 min 峰置信度", 3.12, 5.62, 2.25, GREEN, f"峰体分 {float(a7_2['peak_body_score']):.3f}")
    add_box(slide, 5.56, 5.62, 6.90, 1.14, fill=PALE_GREEN, line="C8E8D2")
    add_text(
        slide,
        "FWHM（Full Width at Half Maximum，半高全宽）是峰高 50% 处左右交点的时间差，用来描述峰宽，不能单独判定毛刺。"
        "本例 FWHM 约 0.071/0.093 min、对称度 0.994/0.973；去尖后峰高仍保留 "
        f"{float(a7_1['despiked_height_retention'])*100:.1f}% / {float(a7_2['despiked_height_retention'])*100:.1f}%。",
        5.82,
        5.89,
        6.40,
        0.58,
        size=10.7,
        bold=True,
    )

    # 13 A3
    slide = new_slide(prs, "代表性工况｜A3：鼓包与背景涟漪被分层处理", 13, source)
    a3_summary = file_near(summary, "A3-")
    a3_plot = find_image(plots, "A3-")
    add_case_banner(slide, "鼓包成果", "宽背景单独标黄；弱密集波动保留为结构化背景残差，不自我确认。", ORANGE)
    if a3_plot:
        add_picture(slide, crop_image(a3_plot, (0.03, 0.02, 0.98, 0.99)), 0.55, 1.68, 8.12, 4.92)
    add_card(slide, f"{float(a3_summary['baseline_excursion_ratio']):.3f}", "基线最大起伏／量程", 8.92, 1.78, 3.18, ORANGE)
    add_card(slide, f"{float(a3_summary['bump_duration']):.2f} min", "最长鼓包时长", 8.92, 3.12, 3.18, ORANGE)
    add_card(slide, str(int(a3_summary["confirmed_peaks"])), "确认峰", 8.92, 4.46, 1.47, GREEN)
    add_card(slide, str(int(a3_summary["structured_background_residuals"])), "背景残差", 10.62, 4.46, 1.47, MUTED)
    add_text(slide, "首尾值接近并不代表无鼓包；A3 由中间最大起伏和持续区间被识别。", 8.95, 5.91, 3.17, 0.57, size=11.8, color=NAVY, bold=True)

    # 14 F1
    slide = new_slide(prs, "代表性工况｜F1：正峰与负峰使用有符号分支共同输出", 14, source)
    f1_neg = event_near(features, "F1-H2", 0.326666667)
    f1_pos = event_near(features, "F1-H2", 0.398333333)
    f1_plot = find_image(plots, "F1-H2")
    add_case_banner(slide, "正负峰", "强负峰达到标准即可确认，不再因为“负号”自动进入 review。", PURPLE)
    if f1_plot:
        add_picture(slide, crop_image(f1_plot, (0.03, 0.02, 0.98, 0.99)), 0.55, 1.70, 8.40, 4.95)
    add_card(slide, f"{float(f1_neg['apex_time_min']):.4f} min", "确认负峰峰顶", 9.18, 1.82, 3.10, PURPLE, f"置信度 {float(f1_neg['peak_confidence_percent']):.1f}%")
    add_card(slide, f"{float(f1_pos['apex_time_min']):.4f} min", "确认正峰峰顶", 9.18, 3.20, 3.10, GREEN, f"置信度 {float(f1_pos['peak_confidence_percent']):.1f}%")
    add_box(slide, 9.18, 4.60, 3.10, 1.64, fill=PALE)
    add_lines(
        slide,
        [
            ("校正值保留符号", NAVY, True),
            (f"负峰 {float(f1_neg['corrected_apex_value']):.4g}", PURPLE, False),
            (f"正峰 {float(f1_pos['corrected_apex_value']):.4g}", GREEN, False),
            ("确认峰同步标注起止边界。", INK, False),
        ],
        9.39,
        4.82,
        2.70,
        1.22,
        size=11.4,
        spacing=3,
    )

    # 15 H1
    slide = new_slide(prs, "代表性工况｜H1：强电信号干扰保持伪峰状态", 15, source)
    h1_plot = find_image(plots, "H1-")
    add_case_banner(slide, "电干扰", "峰宽、采样支持、顶部能量、斜率跳变、对称性与脉冲上下文联合判断。", RED)
    if h1_plot:
        add_picture(slide, crop_image(h1_plot, (0.03, 0.25, 0.98, 0.99)), 0.55, 1.68, 8.12, 4.88)
    for index, event in enumerate(event_near(features, "H1-", t) for t in (3.933333333, 6.14, 6.533333333)):
        add_card(
            slide,
            f"{float(event['apex_time_min']):.4f}",
            "min · 电尖峰",
            8.92,
            1.76 + index * 1.35,
            3.22,
            RED,
            f"artifact · 峰置信度 {float(event['peak_confidence_percent']):.1f}%",
        )
    add_text(slide, "保留时间命中不能覆盖强电形态证据；干扰仍进入可审计 CSV。", 8.94, 5.94, 3.20, 0.48, size=11.7, color=NAVY, bold=True)

    # 16 H2 ambiguity
    slide = new_slide(prs, "重点歧义｜H2 约 0.600 min：算法结论不等于人工真值", 16, source)
    h2 = event_near(features, "H2-", 0.6000)
    h2_plot = find_image(plots, "H2-")
    h2_zoom = find_image(zoom_root, "H2-")
    add_box(slide, 0.55, 1.00, 12.08, 0.62, fill=PALE_RED, line="F4BCC2")
    add_text(slide, "重要：当前输出是“电干扰候选／待复核”，不是已经确认的人工真值。", 0.80, 1.16, 11.55, 0.26, size=14.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    if h2_zoom:
        add_picture(slide, crop_image(h2_zoom, (0.03, 0.04, 0.98, 0.36)), 0.55, 1.82, 6.05, 3.16)
    if h2_plot:
        add_picture(slide, crop_image(h2_plot, (0.03, 0.25, 0.98, 0.99)), 6.78, 1.82, 5.84, 3.16)
    add_box(slide, 0.55, 5.18, 12.08, 1.52, fill=PALE_YELLOW, line="F0D78A")
    add_lines(
        slide,
        [
            (f"算法输出：{h2['feature_type_cn']} / {h2['status_cn']} / 峰置信度 {float(h2['peak_confidence_percent']):.1f}%", ORANGE, True),
            (f"“像峰”证据：FWHM {float(h2['fwhm_min']):.4f} min、对称度 {float(h2['symmetry']):.3f}、峰体分 {float(h2['peak_body_score']):.3f}。", INK, False),
            (f"“异常”证据：峰顶异常={bool(h2['apex_spike_flag'])}，反事实峰体门={bool(h2['peak_body_survives_despike'])}；电干扰分仅 {float(h2['electrical_interference_score']):.3f}。", INK, False),
            ("正确处置：保留 review，建立人工标签后再校准；不能仅凭文件名或算法类别宣称它一定是干扰。", RED, True),
        ],
        0.81,
        5.38,
        11.52,
        1.14,
        size=11.1,
        spacing=3,
    )

    # 17 Noise/background
    slide = new_slide(prs, "噪声与结构化背景：减少“把背景涟漪数成峰”", 17, source)
    add_picture(slide, type_chart(features), 0.55, 1.08, 6.26, 4.98)
    add_picture(slide, confidence_chart(features), 6.95, 1.08, 5.65, 4.02)
    residuals = int((features["feature_type"] == "structured_background_residual").sum())
    add_box(slide, 7.15, 5.30, 5.26, 1.26, fill=PALE)
    add_lines(
        slide,
        [
            (f"结构化背景残差：{residuals} 个", MUTED, True),
            ("它们显示为灰色审计点，置信度受限，不会在鼓包区内自行升级为确认峰。", INK, False),
            ("高噪声也不是一票否决：强而稳定的模板峰仍可确认。", NAVY, True),
        ],
        7.37,
        5.50,
        4.82,
        0.90,
        size=11.3,
        spacing=4,
    )

    # 18 Raw zoom
    slide = new_slide(prs, "原始分段放大图：验证弱峰与局部波动，不改变任何数据", 18, source)
    b10_zoom = find_image(zoom_root, "B10-")
    if b10_zoom:
        add_picture(slide, b10_zoom, 0.55, 1.07, 8.22, 5.78)
        zoom_note = "正式输出中已检测到 41 张 raw_zoom_plots。"
    else:
        add_box(slide, 0.55, 1.07, 8.22, 5.78, fill=PALE)
        add_text(slide, "当前未检测到 raw_zoom_plots；重新运行最新检测器后会自动嵌入。", 1.20, 3.35, 6.90, 0.65, size=16, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        zoom_note = "脚本已实现运行时缺图检测。"
    add_box(slide, 9.02, 1.07, 3.55, 5.78, fill=PALE_BLUE, line="C9DDF7")
    add_lines(
        slide,
        [
            ("用途", BLUE, True),
            ("观察全量程图中被强峰压扁的弱信号、噪声、台阶和局部漂移。", INK, False),
            ("数据原则", BLUE, True),
            ("只缩放显示坐标轴；仍是原始 CSV 样本，不平滑、不插值、不扣基线。", INK, False),
            ("注意", ORANGE, True),
            ("各分段纵轴独立，不能直接比较屏幕高度；定量应回到 CSV。", INK, False),
            ("自动化", BLUE, True),
            (zoom_note, INK, False),
        ],
        9.27,
        1.36,
        3.02,
        5.10,
        size=12.0,
        spacing=6,
    )

    # 19 Final result comparison by workload
    slide = new_slide(prs, "最终检测结果横向对比｜五类工况的输出构成", 19, source)
    add_picture(slide, result_mix_by_workload_chart(features), 0.55, 1.10, 7.75, 5.45)
    workload_counts = {
        folder: Counter(features.loc[features["folder"] == folder, "status"].astype(str))
        for folder in (
            "1-基准谱图类型",
            "2-鼓包谱图类型",
            "3-噪音谱图类型",
            "4-负峰类型",
            "5-电信号干扰类型",
        )
    }
    ref_mix = workload_counts["1-基准谱图类型"]
    bump_mix = workload_counts["2-鼓包谱图类型"]
    noise_mix = workload_counts["3-噪音谱图类型"]
    negative_mix = workload_counts["4-负峰类型"]
    electrical_mix = workload_counts["5-电信号干扰类型"]
    add_box(slide, 8.55, 1.10, 4.03, 5.45, fill=PALE_BLUE, line="C9DDF7")
    add_lines(
        slide,
        [
            ("如何解读", BLUE, True),
            (f"参考标气：{ref_mix['confirmed']}/{sum(ref_mix.values())} 个候选被确认，未输出 artifact。", INK, False),
            (f"鼓包：{bump_mix['review']}/{sum(bump_mix.values())} 进入复核；宽背景被保留而非强行算峰。", INK, False),
            (f"噪声：候选最多（n={sum(noise_mix.values())}），review 占主导，体现保守策略。", INK, False),
            (f"负峰：确认 {negative_mix['confirmed']} 个；负号本身不再触发降级。", INK, False),
            (f"电干扰：artifact {electrical_mix['artifact']} 个，占比为五类中最高。", INK, False),
            ("重要边界", RED, True),
            ("这是算法最终输出构成，不是准确率、召回率或人工真值对比；建立逐峰标签后才能计算性能指标。", RED, True),
        ],
        8.80,
        1.38,
        3.52,
        4.86,
        size=11.5,
        spacing=6,
    )

    # 20 A7 versus H1 comparison
    slide = new_slide(prs, "最终检测对比｜相似的尖顶，为什么 A7 是峰而 H1 是干扰", 20, source)
    add_comparison_headers(slide)
    h1_zoom = find_image(zoom_root, "H1-")
    h1_artifacts = int(
        ((features["file"].astype(str).str.startswith("H1-")) & (features["status"] == "artifact")).sum()
    )
    add_comparison_row(
        slide,
        y=1.30,
        case_label="A7",
        raw_image=a7_zoom,
        raw_crop=(0.03, 0.07, 0.98, 0.34),
        detected_image=a7_plot,
        detected_crop=(0.03, 0.25, 0.98, 0.99),
        conclusion=[
            ("连续峰体", GREEN, True),
            ("0.4333 / 0.5650 min", INK, False),
            (f"置信度 {float(a7_1['peak_confidence_percent']):.1f}% / {float(a7_2['peak_confidence_percent']):.1f}%", INK, False),
            ("去尖后峰体仍在", INK, False),
            ("→ 均确认", GREEN, True),
        ],
        color=GREEN,
    )
    add_comparison_row(
        slide,
        y=4.02,
        case_label="H1",
        raw_image=h1_zoom,
        raw_crop=(0.03, 0.24, 0.98, 0.65),
        detected_image=h1_plot,
        detected_crop=(0.03, 0.25, 0.98, 0.99),
        conclusion=[
            ("少点脉冲", RED, True),
            ("斜率突变且顶部能量集中", INK, False),
            (f"artifact 共 {h1_artifacts} 个", INK, False),
            ("保留时间不能覆盖强异常证据", INK, False),
            ("→ 排除/复核", RED, True),
        ],
        color=RED,
    )
    add_text(
        slide,
        "对比核心：不能只看“尖不尖”；应检查完整峰体、宽度、连续性、对称性、斜率及去尖后的峰体保留。",
        0.72,
        6.62,
        11.72,
        0.30,
        size=12.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 21 A3 and F1 raw-to-final comparisons
    slide = new_slide(prs, "最终检测对比｜鼓包背景与正负峰从原始数据到最终标注", 21, source)
    add_comparison_headers(slide)
    a3_zoom = find_image(zoom_root, "A3-")
    f1_zoom = find_image(zoom_root, "F1-H2")
    add_comparison_row(
        slide,
        y=1.30,
        case_label="A3",
        raw_image=a3_zoom,
        raw_crop=(0.03, 0.22, 0.98, 0.66),
        detected_image=a3_plot,
        detected_crop=(0.03, 0.25, 0.98, 0.99),
        conclusion=[
            ("明显鼓包", ORANGE, True),
            (f"最大起伏 {float(a3_summary['baseline_excursion_ratio'])*100:.1f}%量程", INK, False),
            (f"持续 {float(a3_summary['bump_duration']):.2f} min", INK, False),
            (f"确认 {int(a3_summary['confirmed_peaks'])} / 背景残差 {int(a3_summary['structured_background_residuals'])}", INK, False),
            ("→ 背景不冒充峰", ORANGE, True),
        ],
        color=ORANGE,
    )
    add_comparison_row(
        slide,
        y=4.02,
        case_label="F1",
        raw_image=f1_zoom,
        raw_crop=(0.03, 0.07, 0.98, 0.35),
        detected_image=f1_plot,
        detected_crop=(0.03, 0.25, 0.98, 0.99),
        conclusion=[
            ("有符号双分支", PURPLE, True),
            (f"负峰 {float(f1_neg['apex_time_min']):.4f} min", INK, False),
            (f"正峰 {float(f1_pos['apex_time_min']):.4f} min", INK, False),
            (f"置信度 {float(f1_neg['peak_confidence_percent']):.1f}% / {float(f1_pos['peak_confidence_percent']):.1f}%", INK, False),
            ("→ 正负均确认", PURPLE, True),
        ],
        color=PURPLE,
    )
    add_text(
        slide,
        "左列始终为原始 CSV 样本的显示轴放大；右列才是经过基线、噪声、候选检测和分类后的最终结果。",
        0.72,
        6.62,
        11.72,
        0.30,
        size=11.8,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 22 Boundary
    slide = new_slide(prs, "当前能力边界与验证责任", 22, source)
    add_box(slide, 0.58, 1.15, 5.90, 5.43, fill=PALE_GREEN, line="C8E8D2")
    add_text(slide, "已经具备", 0.88, 1.42, 2.2, 0.34, size=19, color=GREEN, bold=True)
    add_lines(
        slide,
        [
            ("✓ 动态噪声、双基线与有符号正负分支", INK, False),
            ("✓ 普通峰、窄峰、宽峰、峰上峰、负峰候选", INK, False),
            ("✓ 鼓包区域与结构化背景残差分离", INK, False),
            ("✓ 多参数电干扰与“峰体 + 峰顶尖点”区分", INK, False),
            ("✓ 峰顶、边界、置信度、证据分项和原始值输出", INK, False),
            ("✓ 分析图、原始分段图和批量 CSV 交付", INK, False),
        ],
        0.90,
        1.98,
        5.25,
        3.96,
        size=14.0,
        spacing=10,
    )
    add_box(slide, 6.78, 1.15, 5.84, 5.43, fill=PALE_RED, line="F4BCC2")
    add_text(slide, "仍需真值闭环", 7.08, 1.42, 2.8, 0.34, size=19, color=RED, bold=True)
    add_lines(
        slide,
        [
            ("! T1–T6 尚未绑定不同机型／通道的化学组分", INK, False),
            ("! 普通重叠峰的解卷积与面积分配尚未完成", INK, False),
            ("! 正负重叠后的定量面积尚未完成", INK, False),
            ("! 鼓包与基线报警仍需人工标签验证召回率／误报率", INK, False),
            ("! 电干扰排除率、峰型准确率需要逐峰人工真值", INK, False),
            ("! 置信度是概率样评分，不是统计意义上的“真实概率”", RED, True),
        ],
        7.10,
        1.98,
        5.18,
        3.96,
        size=13.7,
        spacing=10,
    )

    # 23 Deliverables
    slide = new_slide(prs, "交付物与下一步", 23, source)
    add_box(slide, 0.58, 1.12, 7.23, 5.54, fill=PALE)
    add_text(slide, "当前交付包", 0.88, 1.40, 2.5, 0.34, size=19, color=NAVY, bold=True)
    add_lines(
        slide,
        [
            ("源代码与配置：src/、configs/、run_detector.py", INK, False),
            ("全量结果：all_detected_features.csv、confirmed_peaks.csv", INK, False),
            ("复核清单：review_required.csv、interference_candidates.csv", INK, False),
            ("背景审计：structured_background_residuals.csv", INK, False),
            ("图像：plots/ 与 raw_zoom_plots/", INK, False),
            ("说明：README、算法说明、技术总结与本演示稿", INK, False),
            ("重生成：python scripts/generate_deliverables.py", TEAL, True),
        ],
        0.90,
        1.98,
        6.55,
        4.00,
        size=13.7,
        spacing=9,
    )
    add_box(slide, 8.10, 1.12, 4.52, 5.54, fill=PALE_BLUE, line="C9DDF7")
    add_text(slide, "建议下一步", 8.40, 1.40, 2.4, 0.34, size=19, color=BLUE, bold=True)
    add_lines(
        slide,
        [
            ("1  建立逐峰人工真值表", BLUE, True),
            ("峰位、峰型、边界、是否干扰。", INK, False),
            ("2  按机型／通道校准", BLUE, True),
            ("绑定组分名称并验证保留时间漂移。", INK, False),
            ("3  建立量化指标", BLUE, True),
            ("峰召回率、伪峰排除率、边界误差、面积误差。", INK, False),
            ("4  再进入解卷积与面积定量", BLUE, True),
            ("先完成真值闭环，再扩大自动确认范围。", RED, True),
        ],
        8.42,
        1.98,
        3.88,
        4.05,
        size=12.7,
        spacing=6,
    )

    # 24-64 Complete one-curve-per-slide appendix
    if zoom_root is None:
        raise RuntimeError("全量结果附录需要 raw_zoom_plots，但当前未找到该目录")
    appendix_rows = summary.sort_values(["folder", "file"], kind="stable").reset_index(drop=True)
    if len(appendix_rows) != 41 or appendix_rows["sha256"].nunique() != 41:
        raise RuntimeError(
            f"全量结果附录要求 41 条独立曲线，实际 rows={len(appendix_rows)}, "
            f"unique_sha={appendix_rows['sha256'].nunique()}"
        )
    for appendix_index, row in appendix_rows.iterrows():
        sha8 = str(row["sha256"])[:8]
        plot_matches = sorted(plots.glob(f"*{sha8}*.png"))
        zoom_matches = sorted(zoom_root.glob(f"*{sha8}*.png"))
        if len(plot_matches) != 1 or len(zoom_matches) != 1:
            raise RuntimeError(
                f"附录图片无法一一配对：file={row['file']}, sha8={sha8}, "
                f"plots={len(plot_matches)}, raw_zoom={len(zoom_matches)}"
            )
        display_index = appendix_index + 1
        slide_number = 23 + display_index
        title = (
            f"全量结果 {display_index:02d}/41｜"
            f"{row['folder']} / {row['file']}"
        )
        slide = new_slide(
            prs,
            title,
            slide_number,
            source,
            title_size=17.5,
        )
        add_text(
            slide,
            "最终三层分析结果",
            0.55,
            0.96,
            2.25,
            0.24,
            size=12.5,
            color=TEAL,
            bold=True,
        )
        add_text(
            slide,
            (
                f"确认 {int(row['confirmed_peaks'])} · "
                f"复核 {int(row['review_features'])} · "
                f"干扰 {int(row['artifacts'])}"
            ),
            2.62,
            0.98,
            3.95,
            0.20,
            size=9.3,
            color=MUTED,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            "原始数据四段局部放大（未预处理）",
            6.76,
            0.96,
            3.35,
            0.24,
            size=12.5,
            color=BLUE,
            bold=True,
        )
        condition_parts = [
            "鼓包报警" if bool(row["bump_flag"]) else "无鼓包报警",
            "高噪声" if bool(row["high_noise_flag"]) else "非高噪声",
        ]
        add_text(
            slide,
            " · ".join(condition_parts),
            10.08,
            0.98,
            2.72,
            0.20,
            size=9.3,
            color=MUTED,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        add_picture(slide, plot_matches[0], 0.50, 1.25, 6.15, 5.62)
        add_picture(slide, zoom_matches[0], 6.70, 1.25, 6.15, 5.62)
    return prs


def style_docx(document) -> None:
    from docx.oxml.ns import qn
    from docx.shared import Cm

    section = document.sections[0]
    section.top_margin = Cm(2.1)
    section.bottom_margin = Cm(1.9)
    section.left_margin = Cm(2.35)
    section.right_margin = Cm(2.15)
    settings = [
        ("Normal", 10.5, False, INK),
        ("Title", 24, True, NAVY),
        ("Heading 1", 18, True, NAVY),
        ("Heading 2", 14, True, TEAL),
        ("Heading 3", 12, True, NAVY2),
    ]
    for name, size, bold, color in settings:
        style = document.styles[name]
        style.font.name = FONT
        style.font.size = Pt(size)
        style.font.bold = bold
        style.font.color.rgb = docx_rgb(color)
        style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)


def markdown_runs(paragraph, text: str) -> None:
    from docx.oxml.ns import qn

    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    for part in re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.bold = True
        elif part.startswith("`") and part.endswith("`"):
            run.text = part[1:-1]
            run.font.name = "Consolas"
            run.font.color.rgb = docx_rgb(PURPLE)
        else:
            run.text = part
            run.font.name = FONT
            run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)


def markdown_to_docx(source: Path, destination: Path) -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm

    document = Document()
    style_docx(document)
    header = document.sections[0].header.paragraphs[0]
    header.text = "ChromPeak 色谱峰检测技术总结 · v0.5.2"
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    for run in header.runs:
        run.font.name = FONT
        run.font.size = Pt(8.5)
        run.font.color.rgb = docx_rgb(MUTED)
    footer = document.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run("第 ")
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instruction = OxmlElement("w:instrText")
    instruction.set(qn("xml:space"), "preserve")
    instruction.text = " PAGE "
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.extend((begin, instruction, end))
    footer.add_run(" 页")

    lines = source.read_text(encoding="utf-8").splitlines()
    index = 0
    code: list[str] = []
    in_code = False
    first_title = True
    while index < len(lines):
        line = lines[index].rstrip()
        if line.startswith("```"):
            if in_code:
                paragraph = document.add_paragraph()
                paragraph.paragraph_format.left_indent = Cm(0.6)
                run = paragraph.add_run("\n".join(code))
                run.font.name = "Consolas"
                run.font.size = Pt(8.5)
                run.font.color.rgb = docx_rgb(PURPLE)
                code.clear()
                in_code = False
            else:
                in_code = True
            index += 1
            continue
        if in_code:
            code.append(line)
            index += 1
            continue
        if not line.strip():
            index += 1
            continue
        if line.startswith("|") and index + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-+", lines[index + 1]):
            headers = [cell.strip() for cell in line.strip("|").split("|")]
            index += 2
            rows = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                rows.append([cell.strip() for cell in lines[index].strip().strip("|").split("|")])
                index += 1
            table = document.add_table(rows=1, cols=len(headers))
            table.style = "Light Shading Accent 1"
            for cell, value in zip(table.rows[0].cells, headers):
                cell.text = re.sub(r"[*`]", "", value)
            for values in rows:
                cells = table.add_row().cells
                for cell, value in zip(cells, values):
                    cell.text = re.sub(r"[*`]", "", value)
            document.add_paragraph()
            continue
        heading = re.match(r"^(#{1,3})\s+(.+)$", line)
        if heading:
            level = len(heading.group(1))
            text = heading.group(2)
            if first_title and level == 1:
                paragraph = document.add_paragraph(style="Title")
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                markdown_runs(paragraph, text)
                subtitle = document.add_paragraph()
                subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
                subrun = subtitle.add_run("基于当前 v0.5.2 正式输出排版")
                subrun.font.name = FONT
                subrun.font.size = Pt(10)
                subrun.font.color.rgb = docx_rgb(MUTED)
                first_title = False
            else:
                paragraph = document.add_paragraph(style=f"Heading {level}")
                markdown_runs(paragraph, text)
            index += 1
            continue
        bullet = re.match(r"^\s*[-*]\s+(.+)$", line)
        numbered = re.match(r"^\s*\d+[.)]\s+(.+)$", line)
        if bullet or numbered:
            paragraph = document.add_paragraph(style="List Bullet" if bullet else "List Number")
            markdown_runs(paragraph, (bullet or numbered).group(1))
            index += 1
            continue
        if line.startswith(">"):
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.left_indent = Cm(0.7)
            markdown_runs(paragraph, line.lstrip("> "))
            for run in paragraph.runs:
                run.italic = True
                run.font.color.rgb = docx_rgb(MUTED)
            index += 1
            continue
        block = [line]
        index += 1
        while index < len(lines):
            next_line = lines[index].rstrip()
            if (
                not next_line.strip()
                or next_line.startswith(("#", "```", "|", ">"))
                or re.match(r"^\s*[-*]\s+", next_line)
                or re.match(r"^\s*\d+[.)]\s+", next_line)
            ):
                break
            block.append(next_line)
            index += 1
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.line_spacing = 1.18
        paragraph.paragraph_format.space_after = Pt(6)
        markdown_runs(paragraph, " ".join(block))
    destination.parent.mkdir(parents=True, exist_ok=True)
    document.save(destination)


def validate_pptx(path: Path) -> dict:
    presentation = Presentation(path)
    slides = len(presentation.slides)
    ratio = float(presentation.slide_width) / float(presentation.slide_height)
    if slides != 64:
        raise RuntimeError(f"PPT 页数异常：{slides}")
    if abs(ratio - 16 / 9) > 0.005:
        raise RuntimeError(f"PPT 比例异常：{ratio:.5f}")
    return {"slides": slides, "ratio": ratio, "bytes": path.stat().st_size}


def validate_docx(path: Path) -> dict:
    from docx import Document

    document = Document(path)
    if not document.paragraphs:
        raise RuntimeError("DOCX 内容为空")
    return {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "bytes": path.stat().st_size,
    }


def main() -> int:
    if not MD_PATH.exists():
        raise FileNotFoundError(f"技术总结尚未生成：{MD_PATH}")
    result_root = select_result_root()
    zoom_root = select_zoom_root(result_root)
    presentation = build_presentation(result_root, zoom_root)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(PPTX_PATH)
    markdown_to_docx(MD_PATH, DOCX_PATH)
    pptx_check = validate_pptx(PPTX_PATH)
    docx_check = validate_docx(DOCX_PATH)
    report = {
        "result_root": str(result_root),
        "zoom_root": str(zoom_root) if zoom_root else None,
        "pptx": str(PPTX_PATH),
        "pptx_validation": pptx_check,
        "docx": str(DOCX_PATH),
        "docx_validation": docx_check,
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
